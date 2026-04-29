"""Tests for the PowerPoint .pptx converter (M32)."""

from __future__ import annotations

from datetime import datetime
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from gateway import converters
from gateway import frontmatter as fm
from gateway import paths
from gateway.converters import pptx as pptx_mod
from gateway.converters.base import ConversionError


@pytest.fixture(autouse=True)
def _reset_registry():
    converters.reset_registry_for_tests()
    yield
    converters.reset_registry_for_tests()


# Layout indexes in the default python-pptx template:
#   0 = title slide (title + subtitle)
#   1 = title and content (title + body)
#   5 = title only
#   6 = blank


def _build_pptx(
    path: Path,
    *,
    slides: list[dict],
    creator: str = "",
    title: str = "",
    created: datetime | None = None,
) -> None:
    """Build a .pptx. Each slide dict supports:
       title, body (list[str]), notes (str), table (list[list[str]]), layout (int).
    """
    prs = Presentation()
    for s in slides:
        layout = prs.slide_layouts[s.get("layout", 1)]
        slide = prs.slides.add_slide(layout)
        if s.get("title") and slide.shapes.title is not None:
            slide.shapes.title.text = s["title"]
        if s.get("body"):
            # The body placeholder for layout 1 is at index 1
            placeholders = list(slide.placeholders)
            body_ph = next(
                (p for p in placeholders if p.placeholder_format.idx == 1),
                None,
            )
            if body_ph is not None:
                tf = body_ph.text_frame
                # First paragraph already exists; subsequent need add_paragraph
                first, *rest = s["body"]
                tf.text = first
                for line in rest:
                    p = tf.add_paragraph()
                    p.text = line
        if s.get("table"):
            grid = s["table"]
            rows = len(grid)
            cols = len(grid[0]) if grid else 0
            shape = slide.shapes.add_table(rows, cols, Inches(1), Inches(3),
                                            Inches(8), Inches(2))
            for r_idx, row in enumerate(grid):
                for c_idx, cell in enumerate(row):
                    shape.table.rows[r_idx].cells[c_idx].text = cell
        if s.get("notes"):
            slide.notes_slide.notes_text_frame.text = s["notes"]
    cp = prs.core_properties
    if creator:
        cp.author = creator
    if title:
        cp.title = title
    if created is not None:
        cp.created = created
    prs.save(str(path))


# --- detect ----------------------------------------------------------------


def test_pptx_detect_accepts_pptx_paths():
    c = pptx_mod.PptxConverter()
    assert c.detect("/tmp/foo.pptx")
    assert c.detect("foo.PPTX")


def test_pptx_detect_rejects_other_extensions_and_urls():
    c = pptx_mod.PptxConverter()
    assert not c.detect("/tmp/foo.ppt")  # legacy not supported
    assert not c.detect("/tmp/foo.docx")
    assert not c.detect("https://example.com/foo.pptx")


# --- convert ---------------------------------------------------------------


def test_pptx_convert_extracts_slides_titles_and_body(kb_root, tmp_path):
    src = tmp_path / "deck.pptx"
    _build_pptx(
        src,
        slides=[
            {"title": "Cover", "body": ["A presentation"]},
            {"title": "Agenda", "body": ["Section A", "Section B", "Section C"]},
        ],
        creator="Andrew Grant",
        title="Launch Deck",
        created=datetime(2026, 4, 28),
    )
    text = pptx_mod.PptxConverter().convert(str(src))
    front, body = fm.parse(text)

    assert front["type"] == "pptx"
    assert front["title"] == "Launch Deck"
    assert front["authors"] == ["Andrew Grant"]
    assert front["published_at"] == "2026"
    assert front["meta"]["slide_count"] == 2
    assert front["meta"]["slides_with_notes"] == 0
    assert front["meta"]["extraction_tool"] == "python-pptx"

    assert "## Slide 1: Cover" in body
    assert "## Slide 2: Agenda" in body
    assert "A presentation" in body
    assert "Section B" in body

    sidecar = paths.raw_dir_for("pptx") / f"{front['id']}.pptx"
    assert sidecar.exists()
    assert sidecar.read_bytes() == src.read_bytes()


def test_pptx_convert_does_not_duplicate_title_in_body(kb_root, tmp_path):
    """The slide title shows up in the heading; it must not also appear as body text.

    Regression: python-pptx returns fresh wrapper objects each access, so an
    `is`-based exclusion of the title shape silently failed.
    """
    src = tmp_path / "deck.pptx"
    unique_title = "M32-Unique-Title-XYZ"
    _build_pptx(
        src,
        slides=[{"title": unique_title, "body": ["body text only"]}],
    )
    text = pptx_mod.PptxConverter().convert(str(src))
    _, body = fm.parse(text)

    # Title appears once (in the heading), zero additional times in body
    assert body.count(unique_title) == 1
    assert "body text only" in body


def test_pptx_convert_captures_speaker_notes(kb_root, tmp_path):
    src = tmp_path / "deck.pptx"
    _build_pptx(
        src,
        slides=[
            {"title": "Slide One", "body": ["body line"], "notes": "Remember the demo timing."},
            {"title": "Slide Two", "body": ["second"]},
        ],
    )
    text = pptx_mod.PptxConverter().convert(str(src))
    front, body = fm.parse(text)

    assert front["meta"]["slides_with_notes"] == 1
    assert "> **Speaker notes:** Remember the demo timing." in body


def test_pptx_convert_renders_table_in_slide(kb_root, tmp_path):
    src = tmp_path / "deck.pptx"
    _build_pptx(
        src,
        slides=[{
            "title": "Numbers",
            "table": [["region", "qty"], ["NA", "12"], ["EU", "7"]],
            "layout": 5,
        }],
    )
    text = pptx_mod.PptxConverter().convert(str(src))
    front, body = fm.parse(text)
    assert "## Slide 1: Numbers" in body
    assert "| region | qty |" in body
    assert "| NA | 12 |" in body


def test_pptx_convert_uses_author_year_title_id(kb_root, tmp_path):
    src = tmp_path / "deck.pptx"
    _build_pptx(
        src,
        slides=[{"title": "S", "body": ["x"]}],
        creator="Smith",
        title="The Big Idea",
        created=datetime(2025, 1, 1),
    )
    text = pptx_mod.PptxConverter().convert(str(src))
    front, _ = fm.parse(text)
    assert front["id"].startswith("pptx-smith-2025-")


def test_pptx_convert_falls_back_to_hash_id_without_metadata(kb_root, tmp_path):
    src = tmp_path / "deck.pptx"
    _build_pptx(src, slides=[{"title": "S", "body": ["x"]}])
    text = pptx_mod.PptxConverter().convert(str(src))
    front, _ = fm.parse(text)
    suffix = front["id"][len("pptx-"):]
    assert len(suffix) == 12
    assert all(c in "0123456789abcdef" for c in suffix)


def test_pptx_convert_escapes_pipe_in_table_cell(kb_root, tmp_path):
    src = tmp_path / "deck.pptx"
    _build_pptx(
        src,
        slides=[{
            "title": "Pipes",
            "table": [["name", "note"], ["foo", "a|b|c"]],
            "layout": 5,
        }],
    )
    text = pptx_mod.PptxConverter().convert(str(src))
    _, body = fm.parse(text)
    assert "a\\|b\\|c" in body


def test_pptx_convert_idempotent_id_for_same_content(kb_root, tmp_path):
    """Two byte-identical .pptx files produce the same canonical ID."""
    src1 = tmp_path / "first.pptx"
    src2 = tmp_path / "second.pptx"
    _build_pptx(src1, slides=[{"title": "S", "body": ["x"]}])
    src2.write_bytes(src1.read_bytes())
    f1, _ = fm.parse(pptx_mod.PptxConverter().convert(str(src1)))
    f2, _ = fm.parse(pptx_mod.PptxConverter().convert(str(src2)))
    assert f1["id"] == f2["id"]


def test_pptx_convert_missing_file_raises(kb_root):
    with pytest.raises(ConversionError):
        pptx_mod.PptxConverter().convert("/does/not/exist.pptx")


def test_pptx_convert_corrupt_file_raises(kb_root, tmp_path):
    src = tmp_path / "corrupt.pptx"
    src.write_bytes(b"not a real pptx file at all")
    with pytest.raises(ConversionError):
        pptx_mod.PptxConverter().convert(str(src))


# --- registry --------------------------------------------------------------


def test_dispatch_pptx_path_routes_to_pptx_converter():
    c = converters.dispatch("/tmp/deck.pptx")
    assert isinstance(c, pptx_mod.PptxConverter)


def test_dispatch_ppt_legacy_format_has_no_converter():
    """Legacy .ppt is intentionally unsupported in v1."""
    from gateway.converters import NoConverterError
    with pytest.raises(NoConverterError):
        converters.dispatch("/tmp/old.ppt")
