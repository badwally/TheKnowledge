"""PowerPoint .pptx converter: presentation → canonical markdown per slide.

Uses python-pptx to walk slides in order. Each slide becomes a `## Slide N`
section, with the slide title used when present. Body text from text-frame
shapes is rendered as paragraphs (preserving paragraph-level indentation
as nested bullets). Tables become markdown tables. Speaker notes, when
present, render as a block-quoted Speaker notes line at the end of the
slide section.

Image and chart shapes are skipped — image ingest and VLM description
land in M33.

Canonical ID per WIKI § 6.1: pptx-<author>-<year>-<short> when metadata
allows, else pptx-<sha256-prefix-12> hash fallback (mirrors PDF / DOCX / XLSX).
"""

from __future__ import annotations

import hashlib
import re
import shutil
from datetime import datetime, timezone
from pathlib import Path
from zipfile import BadZipFile

from pptx import Presentation
from pptx.exc import PackageNotFoundError

from gateway import frontmatter as fm
from gateway import paths, validator
from gateway.converters.base import ConversionError, Converter

_NON_ALNUM = re.compile(r"[^a-z0-9]+")


def _now_iso() -> str:
    return datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")


def _slugify_short(value: str, max_words: int = 4) -> str:
    cleaned = _NON_ALNUM.sub("-", value.lower()).strip("-")
    parts = [p for p in cleaned.split("-") if p]
    return "-".join(parts[:max_words])


def _hash_id_from_bytes(data: bytes, max_chars: int = 12) -> str:
    digest = hashlib.sha256(data).hexdigest()[:max_chars]
    return f"pptx-{digest}"


def _build_id(author: str, title: str, year: str, raw_bytes: bytes) -> str:
    if author and year and title:
        author_slug = _slugify_short(author.split(",")[0], max_words=2)
        title_slug = _slugify_short(title, max_words=3)
        if author_slug and title_slug:
            return f"pptx-{author_slug}-{year}-{title_slug}"
    return _hash_id_from_bytes(raw_bytes)


def _md_escape_cell(value) -> str:
    if value is None:
        return ""
    return (
        str(value)
        .replace("\\", "\\\\")
        .replace("|", "\\|")
        .replace("\r", " ")
        .replace("\n", " ")
    )


def _table_to_markdown(table) -> str:
    rows = []
    for row in table.rows:
        rows.append([cell.text or "" for cell in row.cells])
    if not rows:
        return ""
    headers = rows[0]
    body_rows = rows[1:]
    width = len(headers)
    out = ["| " + " | ".join(_md_escape_cell(h) for h in headers) + " |",
           "| " + " | ".join(["---"] * width) + " |"]
    for r in body_rows:
        padded = list(r[:width])
        while len(padded) < width:
            padded.append("")
        out.append("| " + " | ".join(_md_escape_cell(c) for c in padded) + " |")
    return "\n".join(out)


def _render_paragraph(para) -> str:
    """Render a text-frame paragraph as a markdown line (with indent for level)."""
    text = (para.text or "").strip()
    if not text:
        return ""
    level = getattr(para, "level", 0) or 0
    if level <= 0:
        return text
    return ("  " * level) + "- " + text


def _slide_title(slide) -> str:
    """Best-effort slide title from the title placeholder, else empty."""
    try:
        title_shape = slide.shapes.title
    except (AttributeError, KeyError):
        return ""
    if title_shape is None:
        return ""
    text = (title_shape.text or "").strip()
    return text


def _render_slide(slide, slide_num: int) -> tuple[str, bool]:
    """Render one slide as a markdown section. Returns (markdown, has_notes)."""
    title = _slide_title(slide)
    if title:
        heading = f"## Slide {slide_num}: {title}"
    else:
        heading = f"## Slide {slide_num}"

    parts: list[str] = [heading, ""]
    # Identify the title placeholder by shape_id (Python `is` comparison fails —
    # python-pptx returns fresh wrapper objects on each access).
    title_shape_id: int | None = None
    try:
        ts = slide.shapes.title
        if ts is not None:
            title_shape_id = ts.shape_id
    except (AttributeError, KeyError):
        pass

    body_blocks: list[str] = []
    for shape in slide.shapes:
        # Don't re-emit the title placeholder as a body paragraph
        if title_shape_id is not None and shape.shape_id == title_shape_id:
            continue
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            lines = []
            for para in shape.text_frame.paragraphs:
                rendered = _render_paragraph(para)
                if rendered:
                    lines.append(rendered)
            if lines:
                body_blocks.append("\n".join(lines))
        elif getattr(shape, "has_table", False) and shape.has_table:
            md = _table_to_markdown(shape.table)
            if md:
                body_blocks.append(md)

    if body_blocks:
        parts.append("\n\n".join(body_blocks))

    has_notes = False
    if getattr(slide, "has_notes_slide", False) and slide.has_notes_slide:
        notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
        if notes_text:
            has_notes = True
            parts.append("")
            parts.append("> **Speaker notes:** " + notes_text.replace("\n", " "))

    section = "\n".join(parts).rstrip("\n")
    return section, has_notes


class PptxConverter(Converter):
    type_name = "pptx"

    def detect(self, source: str) -> bool:
        if source.startswith(("http://", "https://")):
            return False
        try:
            return Path(source).suffix.lower() == ".pptx"
        except (TypeError, ValueError):
            return False

    def convert(self, source: str) -> str:
        path = Path(source).expanduser().resolve()
        if not path.exists():
            raise ConversionError(f"pptx not found: {path}")

        raw_bytes = path.read_bytes()
        try:
            prs = Presentation(str(path))
        except (PackageNotFoundError, BadZipFile, KeyError, ValueError) as e:
            raise ConversionError(f"python-pptx failed on {path}: {e}") from e

        slide_sections: list[str] = []
        slide_count = 0
        slides_with_notes = 0
        for idx, slide in enumerate(prs.slides, start=1):
            section, has_notes = _render_slide(slide, idx)
            slide_sections.append(section)
            slide_count += 1
            if has_notes:
                slides_with_notes += 1

        if slide_count == 0:
            raise ConversionError(f"no slides in {path}")

        cp = prs.core_properties
        author = (cp.author or "").strip()
        title = (cp.title or path.stem).strip()
        created = cp.created
        year = created.strftime("%Y") if isinstance(created, datetime) else ""

        body_parts = [f"# {title}",
                      "",
                      f"Presentation with **{slide_count}** slide(s)"
                      + (f", {slides_with_notes} with speaker notes." if slides_with_notes else "."),
                      "",
                      *slide_sections]
        body = "\n\n".join(p for p in body_parts if p != "")
        body = body.rstrip("\n") + "\n"

        canonical_id = _build_id(author, cp.title or "", year, raw_bytes)

        sidecar_dir = paths.raw_dir_for("pptx")
        sidecar_dir.mkdir(parents=True, exist_ok=True)
        sidecar_target = sidecar_dir / f"{canonical_id}.pptx"
        if not sidecar_target.exists():
            shutil.copy2(path, sidecar_target)

        front: dict = {
            "id": canonical_id,
            "type": "pptx",
            "title": title,
            "url": "",
            "authors": _split_authors(author),
            "ingested_at": _now_iso(),
            "content_hash": validator.compute_content_hash(body),
            "source_path": str(sidecar_target.relative_to(paths.knowledge_root())),
            "domains": [],
            "nlm_corpus_ids": [],
            "wiki_pages": [],
            "meta": {
                "slide_count": slide_count,
                "slides_with_notes": slides_with_notes,
                "extraction_tool": "python-pptx",
                "original_filename": path.name,
                "subject": (cp.subject or "").strip(),
            },
        }
        if year:
            front["published_at"] = year
        return fm.serialize(front, body)


def _split_authors(value: str) -> list[str]:
    if not value:
        return []
    parts = [p.strip() for p in re.split(r"[,;]| and ", value)]
    return [p for p in parts if p]
